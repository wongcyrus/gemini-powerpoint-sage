"""Visual generation service for Gemini Powerpoint Sage."""

import io
import logging
import os
import re
from typing import Optional

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from google.adk.agents import LlmAgent
from utils.project_rotation import rotate_project

from config.constants import FilePatterns
from utils.agent_utils import run_visual_agent
from services.visual_generator_helpers import (
    cleanup_reduced_image_file,
    apply_slide_notes,
    build_designer_prompt,
    compute_image_placement_inches,
    get_logo_instruction,
    optimize_image_file,
)

logger = logging.getLogger(__name__)

IMAGE_MODEL_FALLBACKS = (
    "gemini-3.1-flash-image",
    "gemini-3-pro-image",
    "gemini-2.5-flash-image",
)


class VisualGenerator:
    """Service for generating enhanced slide visuals."""

    def __init__(
        self,
        designer_agent: LlmAgent,
        output_dir: str,
        skip_generation: bool = False,
        style: str = "Professional",
    ):
        """
        Initialize the visual generator.

        Args:
            designer_agent: Agent for generating slide designs
            output_dir: Directory to save generated visuals
            skip_generation: Whether to skip visual generation
            style: Style/theme for visual generation
        """
        self.designer_agent = designer_agent
        self.output_dir = output_dir
        self.skip_generation = skip_generation
        self.visual_style = style  # This is visual_style from config
        self.previous_image: Optional[Image.Image] = None

        # Ensure output directory exists
        if not skip_generation:
            os.makedirs(output_dir, exist_ok=True)

    async def generate_visual(
        self,
        slide_idx: int,
        slide_image: Image.Image,
        speaker_notes: str,
        retry_errors: bool = False,
        language: str = "en",
    ) -> Optional[bytes]:
        """
        Generate an enhanced visual for a slide.

        Args:
            slide_idx: Slide index (1-based)
            slide_image: Original slide image
            speaker_notes: Generated speaker notes
            retry_errors: Whether to regenerate existing images
            language: Target language locale code (e.g., en, zh-CN, yue-HK)

        Returns:
            Image bytes if generated, None otherwise
        """
        if self.skip_generation:
            logger.info(
                f"Skipping visual generation for Slide {slide_idx} "
                "(skip-visuals active)."
            )
            return None

        img_filename = FilePatterns.REIMAGINED_SLIDE.format(idx=slide_idx)
        img_path = os.path.join(self.output_dir, img_filename)
        img_bytes = self._load_existing_visual_bytes(img_path, slide_idx, img_filename, retry_errors)
        if img_bytes is None:
            img_bytes = await self._generate_visual_candidate(
                slide_idx,
                slide_image,
                speaker_notes,
                language,
            )

        if img_bytes:
            try:
                with open(img_path, "wb") as f:
                    f.write(img_bytes)
                logger.info("Saved reimagined slide to: %s" % img_path)
                self._update_previous_image(img_bytes)
            except Exception as e:
                logger.error("Failed to save generated image: %s" % e)
        else:
            logger.warning("No image generated for Slide %d" % slide_idx)
            self.previous_image = None

        return img_bytes

    def _load_existing_visual_bytes(
        self,
        img_path: str,
        slide_idx: int,
        img_filename: str,
        retry_errors: bool,
    ) -> Optional[bytes]:
        """Load an existing generated image when reuse is allowed."""
        if not os.path.exists(img_path) or retry_errors:
            return None

        logger.info(
            f"Visual already exists for Slide {slide_idx} "
            f"({img_filename}). Skipping generation."
        )
        try:
            with open(img_path, "rb") as f:
                img_bytes = f.read()
            self._update_previous_image(img_bytes)
            return img_bytes
        except Exception as e:
            logger.error(f"Failed to load existing image {img_path}: {e}")
            return None

    async def _generate_visual_candidate(
        self,
        slide_idx: int,
        slide_image: Image.Image,
        speaker_notes: str,
        language: str,
    ) -> Optional[bytes]:
        """Generate a new visual using the preferred image model order."""

        current_project = rotate_project()
        if current_project:
            logger.info("--- Generating Visual for Slide %d (Project: %s) ---" % (slide_idx, current_project))
        else:
            logger.info("--- Generating Visual for Slide %d ---" % slide_idx)

        logo_instruction = self._get_logo_instruction(slide_idx)
        designer_prompt = self._build_designer_prompt(
            speaker_notes,
            logo_instruction,
            language,
        )
        if designer_prompt is None:
            logger.error("_build_designer_prompt returned None; aborting visual generation.")
            return None

        designer_images = [slide_image]
        if self.previous_image:
            designer_images.append(self.previous_image)

        return await self._generate_with_fallback_models(
            slide_idx=slide_idx,
            prompt=designer_prompt,
            images=designer_images,
        )

    def replace_slide_with_visual(
        self,
        prs,
        slide,
        img_path: str,
        speaker_notes: str,
        mode: str = "cover",  # "contain" or "cover"
        target_dpi: int = 150   # controls file size if you downscale
    ) -> bool:
        """
        Replace slide content with generated visual and add notes to the notes section.

        Args:
            prs: PowerPoint Presentation object
            slide: Slide object to modify
            img_path: Path to the generated image
            speaker_notes: Notes to add to the notes section
            mode: "contain" (preserve aspect, may letterbox) or "cover" (fill, crop overflow)
            target_dpi: DPI to use when re-saving image to reduce file size

        Returns:
            True if successful, False otherwise
        """
        try:
            # Remove all shapes on the slide
            for shape in list(slide.shapes):
                sp = shape.element
                sp.getparent().remove(sp)

            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            logger.debug(f"Slide dimensions: {slide_width.inches}\" x {slide_height.inches}\"")

            # Optionally reduce image file size (not dimensions) by re-saving
            # This keeps pixel dimensions but lowers compression quality for smaller .pptx
            reduced_img_path = optimize_image_file(img_path)

            # Compute placement with aspect ratio rules
            try:
                with Image.open(reduced_img_path) as im:
                    img_width_px, img_height_px = im.size
                    left, top, width, height = self._compute_image_placement(
                        img_width_px, img_height_px, slide_width, slide_height, dpi=96, mode=mode
                    )
            except Exception as e:
                logger.debug(f"Placement fallback (full slide) due to error: {e}")
                left, top, width, height = 0, 0, slide_width, slide_height

            logger.info(f"Original image: {img_path}")
            logger.info(f"Reduced image: {reduced_img_path}")

            # Add the picture with final size AT INSERTION TIME to avoid oversized initial shape
            picture = slide.shapes.add_picture(
                reduced_img_path,
                left=left,
                top=top,
                width=width,
                height=height
            )

            logger.debug(f"Picture dimensions set to: {picture.width.inches}\" x {picture.height.inches}\"")
            logger.info(f"Optimized image added with {mode} mode")

            # Clean up the reduced image if it's different from original
            try:
                cleanup_reduced_image_file(img_path, reduced_img_path)
                if reduced_img_path != img_path:
                    logger.debug(f"Cleaned up reduced image: {reduced_img_path}")
            except Exception as e:
                logger.debug(f"Could not clean up reduced image: {e}")

            # Add speaker notes to the notes section as plain text
            apply_slide_notes(slide, speaker_notes)

            logger.info(f"Replaced slide content with visual using {mode} mode.")
            return True

        except Exception as e:
            logger.error(f"Failed to replace slide with visual: {e}")
            return False





    def add_visual_to_presentation(
        self,
        prs: Presentation,
        slide_idx: int,
        img_path: str,
        speaker_notes: str,
        mode: str = "cover",
        target_dpi: int = 150
    ) -> bool:
        """
        Add a generated visual as a new slide in the presentation.

        Args:
            prs: PowerPoint presentation object
            slide_idx: Original slide index
            img_path: Path to the generated image
            speaker_notes: Speaker notes to add to the slide

        Returns:
            True if successful, False otherwise
        """
        try:
            # Find a blank layout
            try:
                blank_layout = prs.slide_layouts[6]  # Usually blank
            except IndexError:
                logger.warning(
                    "Could not find blank slide layout (index 6), "
                    "using first available."
                )
                blank_layout = prs.slide_layouts[0]

            # Add new slide
            new_slide = prs.slides.add_slide(blank_layout)

            # Get slide dimensions for full coverage
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            logger.debug(f"New slide dimensions: {slide_width.inches}\" x {slide_height.inches}\"")

            # Optimize image file size by re-saving with compression
            reduced_img_path = optimize_image_file(img_path)

            # Compute placement with aspect ratio rules
            try:
                with Image.open(reduced_img_path) as im:
                    img_width_px, img_height_px = im.size
                    left, top, width, height = self._compute_image_placement(
                        img_width_px, img_height_px, slide_width, slide_height, dpi=96, mode=mode
                    )
            except Exception as e:
                logger.debug(f"Placement fallback (full slide) due to error: {e}")
                left, top, width, height = 0, 0, slide_width, slide_height

            # Add the optimized image with proper placement
            picture = new_slide.shapes.add_picture(
                reduced_img_path,
                left=left,
                top=top,
                width=width,
                height=height
            )
            
            logger.debug(f"Picture dimensions set to: {picture.width.inches}\" x {picture.height.inches}\"")
            logger.info(f"Optimized image added with {mode} mode")
            
            # Clean up the reduced image if it's different from original
            try:
                cleanup_reduced_image_file(img_path, reduced_img_path)
                if reduced_img_path != img_path:
                    logger.debug(f"Cleaned up reduced image: {reduced_img_path}")
            except Exception as e:
                logger.debug(f"Could not clean up reduced image: {e}")

            # Add speaker notes to the notes section instead of as text on slide
            apply_slide_notes(
                new_slide,
                f"Generated Notes for Slide {slide_idx}:\n{speaker_notes}",
            )

            logger.info(
                f"Added new slide with full-size reimagined image and notes "
                f"for Slide {slide_idx}."
            )
            return True

        except Exception as e:
            logger.error(f"Failed to add reimagined slide to PPTX: {e}")
            return False

    def _get_logo_instruction(self, slide_idx: int) -> str:
        """Get logo instruction based on slide position."""
        return get_logo_instruction(slide_idx)

    def _build_designer_prompt(
        self,
        speaker_notes: str,
        logo_instruction: str,
        language: str = "en"
    ) -> str:
        """Build the prompt for the primary (Gemini) designer agent."""
        return build_designer_prompt(
            speaker_notes,
            logo_instruction,
            language,
            previous_image_present=self.previous_image is not None,
        )

    def _build_model_agent(self, model_name: str) -> LlmAgent:
        """Build a one-off visual agent for a specific image model."""
        safe_suffix = re.sub(r"[^A-Za-z0-9_]", "_", model_name)
        return LlmAgent(
            name=f"slide_designer_{safe_suffix}",
            model=model_name,
            description="Slide designer",
            instruction=self.designer_agent.instruction,
        )

    @staticmethod
    def _is_429_error(exc: Exception) -> bool:
        """Detect quota/rate-limit failures so we only fallback on 429."""
        current: Exception | None = exc
        while current is not None:
            message = str(current)
            if "429" in message or "RESOURCE_EXHAUSTED" in message or "resource exhausted" in message.lower():
                return True
            current = current.__cause__ if isinstance(current.__cause__, Exception) else None
        return False

    async def _generate_with_fallback_models(
        self,
        *,
        slide_idx: int,
        prompt: str,
        images: list[Image.Image],
    ) -> Optional[bytes]:
        """Try the preferred image models in order, only falling back on 429."""
        attempted_models = []

        for model_name in IMAGE_MODEL_FALLBACKS:
            attempted_models.append(model_name)
            agent = self._build_model_agent(model_name)
            try:
                img_bytes = await run_visual_agent(
                    agent,
                    prompt,
                    images=images,
                    raise_on_error=True,
                )
            except Exception as exc:
                if self._is_429_error(exc):
                    logger.warning(
                        "Image model %s hit 429 for Slide %d; trying next model.",
                        model_name,
                        slide_idx,
                    )
                    continue
                logger.error(
                    "Image generation failed for Slide %d on model %s: %s",
                    slide_idx,
                    model_name,
                    exc,
                )
                return None

            if img_bytes:
                if model_name != IMAGE_MODEL_FALLBACKS[0]:
                    logger.info("Slide %d generated with fallback image model %s", slide_idx, model_name)
                return img_bytes

            logger.error(
                "Image model %s returned no bytes for Slide %d; not retrying unless the error is 429.",
                model_name,
                slide_idx,
            )
            return None

        logger.error(
            "All image models exhausted for Slide %d after attempting: %s",
            slide_idx,
            ", ".join(attempted_models),
        )
        return None

    def _update_previous_image(self, img_bytes: bytes) -> None:
        """Update the previous image for style consistency."""
        try:
            self.previous_image = Image.open(io.BytesIO(img_bytes))
        except Exception as e:
            logger.warning(
                f"Could not load generated image for next iteration: {e}"
            )
            self.previous_image = None

    def reset_style_context(self) -> None:
        """Reset the style context (previous image)."""
        self.previous_image = None

    def _compute_image_placement(
        self, 
        img_width_px: int, 
        img_height_px: int, 
        slide_width, 
        slide_height, 
        dpi: int = 96, 
        mode: str = "cover"
    ):
        """
        Compute image placement on slide with proper aspect ratio handling.
        
        Args:
            img_width_px: Image width in pixels
            img_height_px: Image height in pixels
            slide_width: Slide width (Length object)
            slide_height: Slide height (Length object)
            dpi: DPI for pixel to inch conversion
            mode: "contain" (fit within, may letterbox) or "cover" (fill, may crop)
            
        Returns:
            Tuple of (left, top, width, height) as Length objects
        """
        from pptx.util import Inches

        left_in, top_in, width_in, height_in = compute_image_placement_inches(
            img_width_px,
            img_height_px,
            slide_width.inches,
            slide_height.inches,
            dpi=dpi,
            mode=mode,
        )
        left = Inches(left_in)
        top = Inches(top_in)
        width = Inches(width_in)
        height = Inches(height_in)
        logger.debug(
            f"Image placement ({mode}): {left.inches:.2f}\", {top.inches:.2f}\", {width.inches:.2f}\", {height.inches:.2f}\""
        )
        return left, top, width, height



    def cleanup_temp_files(self, prs) -> None:
        """Clean up temporary image files created during slide processing."""
        cleaned_count = 0
        for slide in prs.slides:
            if hasattr(slide, '_temp_image_paths'):
                for temp_path in slide._temp_image_paths:
                    try:
                        if os.path.exists(temp_path):
                            os.unlink(temp_path)
                            cleaned_count += 1
                            logger.debug(f"Cleaned up temporary file: {temp_path}")
                    except Exception as e:
                        logger.warning(f"Could not clean up {temp_path}: {e}")
                # Clear the list
                slide._temp_image_paths = []
        
        if cleaned_count > 0:
            logger.info(f"Cleaned up {cleaned_count} temporary image files")
