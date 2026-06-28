"""Presentation processing service for Gemini Powerpoint Sage."""

import asyncio
import logging
import os
from typing import Optional, Dict, Any

import pymupdf
from PIL import Image
from pptx import Presentation
from google.genai import types
from google.adk.runners import InMemoryRunner
from google.adk.agents import LlmAgent
from google.adk.tools.agent_tool import AgentTool

from config import Config
from utils.agent_utils import run_stateless_agent, run_visual_agent
from utils.image_utils import register_image, unregister_image
from utils.progress_utils import (
    load_progress,
    save_progress,
    create_slide_key,
    get_progress_file_path,
    should_retry_errors,
)
from utils.project_rotation import rotate_project, get_project_count
from tools.agent_tools import AgentToolFactory
from services.visual_generator import VisualGenerator
from services.tts.tts_orchestrator import TTSOrchestrator, create_tts_orchestrator
from config.tts_config import get_tts_config
from services.presentation_processor_helpers import (
    build_global_context_generation_prompt,
    build_global_context_translation_prompt,
    process_slide_visual,
    build_supervisor_prompt,
    extract_artifact_id,
    extract_video_prompt,
    get_language_name,
    is_error_response,
)
from services.presentation_processor_context_helpers import get_global_context
from services.presentation_processor_output_helpers import save_processed_presentations
from services.presentation_processor_tts_helpers import build_tts_slide_data
import zipfile
import tempfile
import xml.etree.ElementTree as ET

logger = logging.getLogger(__name__)


class PresentationProcessor:
    """Main service for processing presentations and generating speaker notes."""

    def __init__(
        self,
        config: Config,
        supervisor_agent: LlmAgent,
        analyst_agent: LlmAgent,
        writer_agent: LlmAgent,
        auditor_agent: LlmAgent,
        overviewer_agent: LlmAgent,
        designer_agent: LlmAgent,
        translator_agent: Optional[LlmAgent] = None,
        image_translator_agent: Optional[LlmAgent] = None,
        video_generator_agent: Optional[LlmAgent] = None,
    ):
        """
        Initialize the presentation processor.

        Args:
            config: Configuration object
            supervisor_agent: Supervisor agent for orchestration
            analyst_agent: Agent for slide analysis
            writer_agent: Agent for writing speaker notes
            auditor_agent: Agent for auditing existing notes
            overviewer_agent: Agent for generating global context
            designer_agent: Agent for generating visuals
            translator_agent: Agent for translating speaker notes
            image_translator_agent: Agent for translating slide visuals
            video_generator_agent: Agent for generating promotional videos
        """
        self.config = config
        self.supervisor_agent = supervisor_agent
        self.analyst_agent = analyst_agent
        self.writer_agent = writer_agent
        self.auditor_agent = auditor_agent
        self.overviewer_agent = overviewer_agent
        self.designer_agent = designer_agent
        self.translator_agent = translator_agent
        self.image_translator_agent = image_translator_agent
        self.video_generator_agent = video_generator_agent

        # Initialize tool factory
        self.tool_factory = AgentToolFactory(
            analyst_agent=analyst_agent,
            writer_agent=writer_agent,
            auditor_agent=auditor_agent,
            translator_agent=translator_agent,
            image_translator_agent=image_translator_agent,
        )

        # Initialize visual generator
        fallback_model = os.getenv("FALLBACK_IMAGEN_MODEL", "imagen-4.0-generate-001")
        self.visual_generator = VisualGenerator(
            designer_agent=designer_agent,
            output_dir=config.visuals_dir,
            skip_generation=config.skip_visuals,
            fallback_imagen_model=fallback_model,
            style=config.visual_style,
        )

        # Initialize TTS orchestrator
        try:
            tts_config = get_tts_config()
            if tts_config.enabled:
                self.tts_orchestrator = create_tts_orchestrator(tts_config, main_config=self.config)
                logger.info("TTS orchestrator initialized successfully")
            else:
                self.tts_orchestrator = None
                logger.info("TTS system is disabled")
        except Exception as e:
            logger.warning(f"Failed to initialize TTS orchestrator: {e}")
            self.tts_orchestrator = None

        # Progress tracking
        self.progress_file = get_progress_file_path(
            config.pptx_path,
            config.language,
            config._get_output_dir()
        )
        self.retry_errors = should_retry_errors()

        logger.info(f"Initialized processor with config: {config}")
        logger.info(f"Progress file: {self.progress_file}")
        logger.info(f"Retry errors: {self.retry_errors}")

    async def process(self) -> tuple[str, str]:
        """
        Process the presentation and generate speaker notes.

        Returns:
            Tuple of (notes_only_path, with_visuals_path)
        """
        logger.info(f"Processing PPTX: {self.config.pptx_path}")
        logger.info(f"Region: {os.environ.get('GOOGLE_CLOUD_LOCATION')}")
        logger.info(f"Language: {self.config.language}")

        # Load files - create two separate presentations
        prs_notes = Presentation(self.config.pptx_path)
        prs_visuals = Presentation(self.config.pptx_path)
        
        # Force 16:9 aspect ratio for visuals presentation immediately
        # This ensures image placement calculations are correct
        from pptx.util import Inches
        prs_visuals.slide_width = Inches(10)
        prs_visuals.slide_height = Inches(5.625)

        pdf_doc = pymupdf.open(self.config.pdf_path)
        
        # Validate PDF document
        pdf_page_count = len(pdf_doc)
        pptx_slide_count = len(prs_notes.slides)
        
        logger.info(f"PDF pages: {pdf_page_count}, PPTX slides: {pptx_slide_count}")
        
        if pdf_page_count == 0:
            raise ValueError("PDF document has no pages")
        
        limit = min(pptx_slide_count, pdf_page_count)
        
        if limit != pptx_slide_count:
            logger.warning(
                f"PDF has fewer pages ({pdf_page_count}) than PPTX slides ({pptx_slide_count}). "
                f"Processing only {limit} slides."
            )

        # Load progress
        progress = load_progress(self.progress_file)

        # Load English notes if processing non-English language
        self.english_notes = {}
        if self.config.language != "en":
            from utils.progress_utils import get_progress_file_path
            en_progress_file = get_progress_file_path(
                self.config.pptx_path, "en", self.config._get_output_dir()
            )
            if os.path.exists(en_progress_file):
                en_progress = load_progress(en_progress_file)
                for slide_data in en_progress.get("slides", {}).values():
                    if slide_data.get("status") == "success":
                        slide_idx = slide_data.get("slide_index")
                        note = slide_data.get("note")
                        if slide_idx and note:
                            self.english_notes[slide_idx] = note
                logger.info(
                    f"Loaded {len(self.english_notes)} English notes "
                    "for translation"
                )
            else:
                logger.warning(
                    f"English notes not found at {en_progress_file}. "
                    "Translation mode will fall back to generation mode."
                )

        # Generate or load global context
        global_context = await self._get_global_context(
            pdf_doc, limit, progress
        )

        # Get presentation theme
        presentation_theme = self.config.get_presentation_theme()

        # Setup supervisor tools
        self._configure_supervisor_tools(
            presentation_theme,
            global_context,
            self.english_notes
        )

        # Initialize supervisor runner
        supervisor_runner = await self._initialize_supervisor()

        # PHASE 1: Generate all speaker notes
        slide_data = await self._phase_generate_notes(
            prs_notes, prs_visuals, pdf_doc, limit, progress,
            supervisor_runner, presentation_theme, global_context
        )

        # PHASE 1.5: Generate TTS audio (if enabled)
        await self._phase_generate_tts(slide_data, progress)

        # PHASE 2: Generate all visuals
        missing_visuals_count = await self._phase_generate_visuals(
            prs_visuals, slide_data
        )

        # PHASE 3: Generate videos (if enabled)
        await self._phase_generate_videos(
            slide_data, presentation_theme, global_context
        )

        # Finalize and save
        return self._save_outputs(
            prs_notes, prs_visuals, missing_visuals_count
        )

    async def _phase_generate_notes(
        self,
        prs_notes: Presentation,
        prs_visuals: Presentation,
        pdf_doc,
        limit: int,
        progress: Dict[str, Any],
        supervisor_runner: InMemoryRunner,
        presentation_theme: str,
        global_context: str
    ) -> list:
        """
        Phase 1: Generate all speaker notes.
        
        Returns:
            List of slide data dictionaries for visual processing.
        """
        logger.info("\n" + "="*60)
        logger.info("PHASE 1: Generating speaker notes for all slides")
        logger.info("="*60)
        
        from utils.pptx_utils import get_slide_notes, update_slide_notes
        
        slide_data = []  # Store slide data for visual processing
        previous_slide_summary = "Start of presentation."
        previous_speaker_notes = []  # Track last 3 slides' full speaker notes
        
        user_id = "supervisor_user"
        session_id = "supervisor_session"

        # Collect all image IDs for cleanup at the end
        image_ids_to_cleanup = []

        for i in range(limit):
            slide_idx = i + 1
            
            # Validate slide and PDF page indices
            if i >= len(prs_notes.slides):
                logger.error(f"Slide index {i} exceeds PPTX slide count ({len(prs_notes.slides)})")
                break
                
            if i >= len(pdf_doc):
                logger.error(f"PDF page index {i} exceeds PDF page count ({len(pdf_doc)})")
                break
            
            slide_notes = prs_notes.slides[i]
            slide_visuals = prs_visuals.slides[i]
            pdf_page = pdf_doc[i]

            # Rotate Google Cloud project before processing each slide
            current_project = rotate_project()
            if current_project:
                logger.info(f"--- Processing Notes for Slide {slide_idx} (Project: {current_project}) ---")
            else:
                logger.info(f"--- Processing Notes for Slide {slide_idx} ---")

            # Get existing notes (from notes presentation)
            existing_notes = get_slide_notes(slide_notes)
            skey = create_slide_key(slide_idx, existing_notes)
            entry = progress["slides"].get(skey)

            # Register slide image
            image_id = f"slide_{slide_idx}"
            image_ids_to_cleanup.append(image_id)
            try:
                slide_image = self._extract_slide_image(pdf_page)
                register_image(image_id, slide_image)
            except Exception as e:
                logger.error(f"Failed to extract/register image for slide {slide_idx}: {e}")
                # Create a placeholder image and register it
                placeholder = Image.new('RGB', (800, 600), color='lightgray')
                register_image(image_id, placeholder)
                slide_image = placeholder

            # Generate or retrieve speaker notes
            final_response, status = await self._process_slide_notes(
                slide_idx, image_id, existing_notes,
                previous_slide_summary, presentation_theme,
                global_context, entry, supervisor_runner,
                user_id, session_id, total_slides=limit,
                previous_speaker_notes=previous_speaker_notes
            )

            # Update slide notes in both presentations
            if status == "success":
                # Update notes in notes-only presentation
                update_slide_notes(slide_notes, final_response)
                
                # Update notes in visuals presentation
                update_slide_notes(slide_visuals, final_response)
                
                previous_slide_summary = final_response[:200]
                
                # Track previous speaker notes (keep last 3)
                previous_speaker_notes.append({
                    'slide_idx': slide_idx,
                    'notes': final_response
                })
                if len(previous_speaker_notes) > 3:
                    previous_speaker_notes.pop(0)

            # Update progress
            progress["slides"][skey] = {
                "slide_index": slide_idx,
                "existing_notes_hash": skey.split("_")[-1],
                "original_notes": existing_notes,
                "note": final_response,
                "status": status,
            }
            save_progress(self.progress_file, progress)

            # Store data for visual generation phase
            slide_data.append({
                "slide_idx": slide_idx,
                "slide_visuals": slide_visuals,
                "slide_image": slide_image,
                "speaker_notes": final_response,
                "status": status,
            })
            
        # Cleanup all images at the end to avoid premature deletion
        for image_id in image_ids_to_cleanup:
            unregister_image(image_id)
            
        return slide_data

    async def _phase_generate_tts(
        self,
        slide_data: list,
        progress: Dict[str, Any]
    ) -> None:
        """
        Phase 1.5: Generate TTS audio for all slides.
        
        Args:
            slide_data: List of slide data from notes generation
            progress: Progress tracking dictionary
        """
        if not self.tts_orchestrator:
            logger.info("TTS generation skipped (TTS system disabled)")
            return
        
        logger.info("\n" + "="*60)
        logger.info("PHASE 1.5: Generating TTS audio for all slides")
        logger.info("="*60)
        
        # Extract presentation ID from config
        presentation_id = os.path.splitext(os.path.basename(self.config.pptx_path))[0]
        
        tts_slides = build_tts_slide_data(slide_data, self.config.language, presentation_id)
        
        if not tts_slides:
            logger.warning("No slides with successful notes found for TTS generation")
            return
        
        try:
            # Generate TTS for all slides
            logger.info(f"Generating TTS for {len(tts_slides)} slides in language {self.config.language}")
            
            tts_results = await self.tts_orchestrator.process_single_language_batch(
                tts_slides, self.config.language, presentation_id
            )
            
            # Update progress with TTS results
            successful_tts = 0
            for i, tts_result in enumerate(tts_results):
                slide_idx = tts_slides[i].slide_number
                
                # Find corresponding progress entry
                for skey, slide_progress in progress["slides"].items():
                    if slide_progress.get("slide_index") == slide_idx:
                        # Add TTS information to progress
                        if tts_result.is_valid():
                            slide_progress["audio_file_path"] = tts_result.file_path
                            slide_progress["tts_metadata"] = {
                                "engine_used": tts_result.engine_used.value,
                                "duration_seconds": tts_result.duration_seconds,
                                "cache_key": tts_result.cache_key,
                                "style_prompt": tts_result.style_prompt
                            }
                            successful_tts += 1
                            logger.debug(f"✓ TTS generated for slide {slide_idx}")
                        else:
                            slide_progress["tts_metadata"] = {
                                "error": tts_result.metadata.get("error", "Unknown TTS error")
                            }
                            logger.warning(f"✗ TTS failed for slide {slide_idx}")
                        break
            
            # Save updated progress
            save_progress(self.progress_file, progress)
            
            logger.info(f"TTS generation completed: {successful_tts}/{len(tts_slides)} successful")
            
            # Log TTS statistics
            stats = self.tts_orchestrator.get_orchestrator_stats()
            logger.info(f"TTS Cache: {stats['cache']['total_entries']} entries, "
                       f"{stats['cache']['total_size_mb']:.1f} MB")
            
        except Exception as e:
            logger.error(f"TTS generation failed: {e}")
            # Continue processing even if TTS fails (graceful degradation)

    async def _phase_generate_visuals(
        self,
        prs_visuals: Presentation,
        slide_data: list
    ) -> int:
        """
        Phase 2: Generate all visuals.
        
        Returns:
            Number of slides with missing visuals.
        """
        logger.info("\n" + "="*60)
        logger.info("PHASE 2: Generating visuals for all slides")
        logger.info("="*60)
        
        missing_visuals_count = 0
        for slide_info in slide_data:
            slide_idx = slide_info["slide_idx"]
            slide_visuals = slide_info["slide_visuals"]
            slide_image = slide_info["slide_image"]
            speaker_notes = slide_info["speaker_notes"]
            status = slide_info["status"]

            # Rotate Google Cloud project before processing each slide visual
            current_project = rotate_project()
            if current_project:
                logger.debug(f"Processing visual for Slide {slide_idx} (Project: {current_project})")

            if status != "success":
                logger.warning(
                    f"Skipping visual generation for Slide {slide_idx} "
                    f"due to notes generation failure"
                )
                missing_visuals_count += 1
                continue

            missing_visuals_count += await process_slide_visual(
                slide_idx=slide_idx,
                slide_visuals=slide_visuals,
                slide_image=slide_image,
                speaker_notes=speaker_notes,
                status=status,
                language=self.config.language,
                visuals_dir=os.path.dirname(self.config.visuals_dir),
                pptx_path=self.config.pptx_path,
                retry_errors=self.retry_errors,
                image_translator_agent=self.image_translator_agent,
                visual_generator=self.visual_generator,
                replace_visual=lambda sv, img_path, notes: self.visual_generator.replace_slide_with_visual(
                    prs_visuals, sv, img_path, notes
                ),
                run_visual_agent=run_visual_agent,
                get_language_name=get_language_name,
            )

        return missing_visuals_count

    async def _phase_generate_videos(
        self,
        slide_data: list,
        presentation_theme: str,
        global_context: str
    ) -> None:
        """
        Phase 3: Generate videos (if enabled).
        """
        if self.config.generate_videos and self.video_generator_agent:
            logger.info("\n" + "="*60)
            logger.info("PHASE 3: Generating videos for all slides")
            logger.info("="*60)
            try:
                await self._generate_videos_for_slides(
                    slide_data, presentation_theme, global_context
                )
            except Exception:
                logger.error("Error during video generation phase", exc_info=True)
                logger.warning("Continuing without videos")

    def _save_outputs(
        self,
        prs_notes: Presentation,
        prs_visuals: Presentation,
        missing_visuals_count: int
    ) -> tuple[str, str]:
        """
        Save the processed presentations to disk.
        """
        from utils.pptx_utils import ensure_pptx_path, restore_vba_project

        output_path_notes, output_path_visuals = save_processed_presentations(
            prs_notes=prs_notes,
            prs_visuals=prs_visuals,
            output_path_notes=self.config.output_path,
            output_path_visuals=self.config.output_path_with_visuals,
            source_pptx_path=self.config.pptx_path,
            missing_visuals_count=missing_visuals_count,
            ensure_pptx_path=ensure_pptx_path,
            restore_vba_project=restore_vba_project,
        )

        logger.info("Saved presentation with notes to: %s", output_path_notes)
        if output_path_visuals:
            logger.info("Saved presentation with visuals to: %s", output_path_visuals)
        else:
            logger.warning(
                "Skipping visuals presentation save: %d slide(s) missing images",
                missing_visuals_count,
            )

        return output_path_notes, output_path_visuals

    async def _get_global_context(
        self,
        pdf_doc,
        limit: int,
        progress: Dict[str, Any]
    ) -> str:
        """Generate or retrieve cached global context."""
        from config.constants import LanguageConfig
        from utils.progress_utils import get_progress_file_path

        return await get_global_context(
            pdf_doc=pdf_doc,
            limit=limit,
            progress=progress,
            language=self.config.language,
            retry_errors=self.retry_errors,
            progress_file=self.progress_file,
            output_dir=self.config._get_output_dir(),
            pptx_path=self.config.pptx_path,
            load_progress=load_progress,
            save_progress=save_progress,
            get_progress_file_path=get_progress_file_path,
            run_stateless_agent=run_stateless_agent,
            overviewer_agent=self.overviewer_agent,
            translator_agent=self.translator_agent,
            build_generation_prompt=build_global_context_generation_prompt,
            build_translation_prompt=build_global_context_translation_prompt,
            language_name_lookup=LanguageConfig.get_language_name,
        )

    def _configure_supervisor_tools(
        self,
        presentation_theme: str,
        global_context: str,
        english_notes: dict = None,
    ) -> None:
        """Configure the supervisor agent's tools."""
        tools = [
            self.tool_factory.create_analyst_tool(),
            self.tool_factory.create_writer_tool(
                presentation_theme,
                global_context,
                self.config.language,
                english_notes,
                self.config.speaker_style
            ),
            self.tool_factory.create_auditor_tool(self.config.language),
        ]
        
        # Add translator tool if available (used for language correction when auditor fails)
        if self.translator_agent:
            tools.append(self.tool_factory.create_translator_tool())
        
        self.supervisor_agent.tools = tools

    async def _initialize_supervisor(self) -> InMemoryRunner:
        """Initialize and create supervisor session."""
        supervisor_runner = InMemoryRunner(
            agent=self.supervisor_agent,
            app_name="agents"
        )

        user_id = "supervisor_user"
        session_id = "supervisor_session"

        await supervisor_runner.session_service.create_session(
            app_name="agents",
            user_id=user_id,
            session_id=session_id
        )

        return supervisor_runner

    def _extract_slide_image(self, pdf_page) -> Image.Image:
        """Extract image from PDF page."""
        try:
            pix = pdf_page.get_pixmap(dpi=150)
            if not pix or pix.width == 0 or pix.height == 0:
                raise ValueError(f"Invalid pixmap dimensions: {pix.width}x{pix.height}")
            
            image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            if not image:
                raise ValueError("Failed to create PIL Image from pixmap data")
                
            return image
        except Exception as e:
            logger.error(f"Failed to extract image from PDF page: {e}")
            # Create a placeholder image to prevent complete failure
            placeholder = Image.new('RGB', (800, 600), color='lightgray')
            return placeholder

    async def _process_slide_notes(
        self,
        slide_idx: int,
        image_id: str,
        existing_notes: str,
        previous_slide_summary: str,
        presentation_theme: str,
        global_context: str,
        entry: Optional[Dict[str, Any]],
        supervisor_runner: InMemoryRunner,
        user_id: str,
        session_id: str,
        total_slides: int = None,
        previous_speaker_notes: list = None,
    ) -> tuple[str, str]:
        """
        Process notes for a single slide.

        Returns:
            Tuple of (final_response, status)
        """
        # TRANSLATION MODE: If non-English and English notes exist
        # translate directly (BEFORE checking if already done)
        if (
            self.config.language != "en"
            and self.translator_agent
            and hasattr(self, 'english_notes')
            and self.english_notes
            and slide_idx in self.english_notes
        ):
            # Check if already translated
            if (
                entry
                and entry.get("status") == "success"
                and not self.retry_errors
            ):
                existing_note = entry.get("note", "")
                # Verify it's actually translated (not English)
                # Simple check: if it matches English note, retranslate
                if existing_note != self.english_notes[slide_idx]:
                    logger.info(
                        f"Skipping translation for slide {slide_idx} "
                        "(already translated)"
                    )
                    return existing_note, "success"

            logger.info(
                f"Translation mode: Translating slide {slide_idx} "
                f"from English to {self.config.language}"
            )
            english_note = self.english_notes[slide_idx]
            translated_note = await self._translate_notes(
                english_note, slide_idx
            )
            if translated_note:
                return translated_note, "success"
            else:
                logger.warning(
                    f"Translation failed for slide {slide_idx}, "
                    "falling back to generation mode"
                )

        # Check if already done (for non-translation mode)
        if (
            entry
            and entry.get("status") == "success"
            and not self.retry_errors
        ):
            logger.info(f"Skipping generation for slide {slide_idx}")
            return entry.get("note", ""), "success"

        # Build supervisor prompt
        supervisor_prompt = self._build_supervisor_prompt(
            slide_idx, image_id, existing_notes,
            previous_slide_summary, presentation_theme,
            global_context, total_slides, previous_speaker_notes
        )

        content = types.Content(
            role='user',
            parts=[types.Part.from_text(text=supervisor_prompt)]
        )

        # Run supervisor with retry logic
        final_response = ""
        status = "pending"
        max_retries = 3
        retry_delay = 2  # seconds

        for attempt in range(max_retries):
            try:
                final_response = ""
                async for event in supervisor_runner.run_async(
                    user_id=user_id,
                    session_id=session_id,
                    new_message=content,
                ):
                    if getattr(event, "content", None) and event.content.parts:
                        for part in event.content.parts:
                            fn_call = getattr(part, "function_call", None)
                            if fn_call:
                                print(
                                    f"\n[Supervisor] 📞 calling tool: "
                                    f"{fn_call.name}"
                                )
                            text = getattr(part, "text", "") or ""
                            final_response += text

                # Check if we got a response
                final_response = final_response.strip()
                if final_response:
                    # Check if the response is an error message from tools
                    if self._is_error_response(final_response):
                        logger.error(
                            f"Tool error detected in response for Slide {slide_idx}: "
                            f"{final_response[:100]}..."
                        )
                        status = "error"
                        break
                    else:
                        status = "success"
                        self.tool_factory.reset_writer_output()
                        break

                # Try fallback to last writer output
                last_output = self.tool_factory.last_writer_output
                if last_output:
                    # Check if fallback output is also an error
                    if self._is_error_response(last_output):
                        logger.error(
                            f"Fallback output is also an error for Slide {slide_idx}: "
                            f"{last_output[:100]}..."
                        )
                        status = "error"
                        break
                    else:
                        logger.info(
                            f"Supervisor returned empty text, "
                            f"using fallback content ({len(last_output)} chars)."
                        )
                        final_response = last_output
                        status = "success"
                        self.tool_factory.reset_writer_output()
                        break

                # No response and no fallback - retry if attempts remain
                if attempt < max_retries - 1:
                    wait_time = retry_delay * (2 ** attempt)
                    logger.warning(
                        f"Empty response for Slide {slide_idx}, "
                        f"retrying in {wait_time}s (attempt {attempt + 1}/"
                        f"{max_retries})..."
                    )
                    await asyncio.sleep(wait_time)
                else:
                    logger.error(
                        f"Failed to get response for Slide {slide_idx} "
                        f"after {max_retries} attempts."
                    )
                    status = "error"

            except Exception as e:
                if attempt < max_retries - 1:
                    wait_time = retry_delay * (2 ** attempt)
                    logger.error(
                        f"Error in supervisor loop (attempt {attempt + 1}/"
                        f"{max_retries}): {e}, retrying in {wait_time}s..."
                    )
                    await asyncio.sleep(wait_time)
                else:
                    logger.error(
                        f"Error in supervisor loop after {max_retries} "
                        f"attempts: {e}"
                    )
                    status = "error"
                    break

        return final_response, status

    def _is_error_response(self, response: str) -> bool:
        """
        Check if a response contains error messages from tools.
        
        Uses only highly specific patterns that are very unlikely to appear 
        in legitimate speaker notes to minimize false positives.
        
        Args:
            response: The response text to check
            
        Returns:
            True if the response appears to be an error message
        """
        return is_error_response(response)

    def _build_supervisor_prompt(
        self,
        slide_idx: int,
        image_id: str,
        existing_notes: str,
        previous_slide_summary: str,
        presentation_theme: str,
        global_context: str,
        total_slides: int = None,
        previous_speaker_notes: list = None,
    ) -> str:
        """Build the prompt for the supervisor agent."""
        return build_supervisor_prompt(
            slide_idx=slide_idx,
            image_id=image_id,
            existing_notes=existing_notes,
            previous_slide_summary=previous_slide_summary,
            presentation_theme=presentation_theme,
            global_context=global_context,
            target_language=self.config.language,
            total_slides=total_slides,
            previous_speaker_notes=previous_speaker_notes,
        )

    async def _translate_notes(
        self, english_note: str, slide_idx: int
    ) -> Optional[str]:
        """
        Translate English speaker notes to target language.

        Args:
            english_note: The English speaker note to translate
            slide_idx: The slide index for context

        Returns:
            Translated notes or None if translation fails
        """
        if not self.translator_agent:
            return None

        lang_name = get_language_name(self.config.language)

        prompt = (
            f"Translate the following English speaker notes to {lang_name}. "
            f"Maintain technical accuracy, educational tone, and clarity. "
            f"Preserve formatting and structure.\n\n"
            f"English notes:\n{english_note}\n\n"
            f"IMPORTANT: Provide ONLY the translated speaker notes "
            f"in {lang_name}. Do not include explanations or metadata."
        )

        try:
            result = await run_stateless_agent(
                self.translator_agent, prompt
            )
            if result and result.strip():
                logger.info(
                    f"Successfully translated slide {slide_idx} "
                    f"to {self.config.language}"
                )
                return result.strip()
        except Exception as e:
            logger.error(
                f"Translation error for slide {slide_idx}: {e}",
                exc_info=True
            )

        return None

    async def _translate_visual(
        self,
        english_visual: Image.Image,
        english_note: str,
        translated_note: str,
        slide_idx: int,
    ) -> Optional[bytes]:
        """
        Translate visual from English to target language.

        Uses Image Translator Agent to analyze English visual and
        Designer Agent to regenerate with translated text.

        Args:
            english_visual: The English slide visual image
            english_note: The English speaker notes for context
            translated_note: The translated speaker notes
            slide_idx: The slide index

        Returns:
            Image bytes of translated visual or None if translation fails
        """
        if not self.image_translator_agent:
            return None

        lang_name = get_language_name(self.config.language)

        # Step 1: Analyze English visual and get translation specs
        analysis_prompt = (
            f"You are a visual translator. Analyze this slide image and "
            f"provide ONLY the following information:\n\n"
            f"TEXT TRANSLATIONS (list all text in the image):\n"
            f"English: [text] -> {lang_name}: [translation]\n\n"
            f"VISUAL DESCRIPTION:\n"
            f"Describe the slide layout, colors, and design elements.\n\n"
            f"Context:\n"
            f"English notes: {english_note}\n"
            f"Translated notes: {translated_note}\n\n"
            f"IMPORTANT: Be concise. List only the text translations and "
            f"visual description. No explanations or planning."
        )

        try:
            # Get translation specs from image translator
            translation_spec = await run_stateless_agent(
                self.image_translator_agent,
                analysis_prompt,
                [english_visual]  # Pass as list
            )

            if not translation_spec or not translation_spec.strip():
                logger.warning(
                    f"Image translator returned empty result for "
                    f"slide {slide_idx}"
                )
                return None

            # Step 2: Use designer agent to regenerate visual
            # with translated content
            design_prompt = (
                f"Generate a slide visual in {lang_name} based on these "
                f"specifications:\n\n"
                f"{translation_spec}\n\n"
                f"Translated speaker notes:\n{translated_note}\n\n"
                f"IMPORTANT:\n"
                f"- Use the translated text from the specifications\n"
                f"- Maintain the same layout and design style\n"
                f"- Ensure all text is in {lang_name}\n"
                f"- Keep colors and branding consistent\n"
                f"- Make it professional and high-quality"
            )

            from utils.agent_utils import run_visual_agent
            img_bytes = await run_visual_agent(
                self.designer_agent,
                design_prompt,
                images=[english_visual]
            )

            if img_bytes:
                logger.info(
                    f"Successfully translated visual for slide {slide_idx} "
                    f"to {self.config.language}"
                )
                return img_bytes
            else:
                logger.warning(
                    f"Designer failed to generate translated visual for "
                    f"slide {slide_idx}"
                )
                return None

        except Exception as e:
            logger.error(
                f"Visual translation error for slide {slide_idx}: {e}",
                exc_info=True
            )
            return None

    async def _generate_videos_for_slides(
        self,
        slide_data: list,
        _presentation_theme: str,
        _global_context: str,
    ) -> None:
        """
        Generate videos for all slides using the video generator agent.
        
        Calls the MCP-backed video agent with slide images and speaker notes
        to generate MP4 videos. Falls back to text prompts if agent unavailable.
        
        Args:
            slide_data: List of slide data dictionaries with slide info
            _presentation_theme: Theme/context for the presentation
            _global_context: Global context about the presentation
        """
        # Ensure videos directory exists
        videos_dir = self.config.videos_dir
        os.makedirs(videos_dir, exist_ok=True)
        logger.info("Videos directory: %s", videos_dir)
        
        # Load slide images from the PDF for video generation
        pdf_doc = pymupdf.open(self.config.pdf_path)
        
        # Process each slide for video generation
        for slide_info in slide_data:
            slide_idx = slide_info["slide_idx"]
            speaker_notes = slide_info["speaker_notes"]
            status = slide_info["status"]
            
            if status != "success":
                logger.warning(
                    "Skipping video for Slide %d (status: %s)",
                    slide_idx, status
                )
                continue
            
            try:
                logger.info("Generating video for Slide %d", slide_idx)
                
                # Extract video prompt
                video_prompt = self._extract_video_prompt(
                    slide_idx, speaker_notes
                )
                
                if not video_prompt or not video_prompt.strip():
                    logger.warning(
                        "Failed to generate video prompt for Slide %d",
                        slide_idx
                    )
                    continue
                
                # Try to call MCP-backed video agent with slide image
                video_data = None
                try:
                    # Load slide image from PDF
                    if slide_idx - 1 < len(pdf_doc):
                        pix = pdf_doc[slide_idx - 1].get_pixmap(dpi=75)
                        slide_img = Image.frombytes(
                            "RGB",
                            [pix.width, pix.height],
                            pix.samples
                        )
                        
                        # Call video agent with image and prompt
                        from utils.agent_utils import run_stateless_agent
                        
                        agent_prompt = (
                            f"Generate a professional video for a presentation "
                            f"slide based on this concept:\n\n{video_prompt}\n\n"
                            f"Speaker Notes:\n{speaker_notes}\n\n"
                            f"Use the slide image provided to guide the visual "
                            f"style. Generate an 8-10 second video."
                        )
                        
                        response = await run_stateless_agent(
                            self.video_generator_agent,
                            agent_prompt,
                            images=[slide_img]
                        )
                        
                        logger.info(
                            "Video agent response for Slide %d: %s",
                            slide_idx, response[:200]
                        )
                        
                        # Parse response for artifact or video data
                        # The MCP agent should return artifact_id if successful
                        video_artifact_id = self._extract_artifact_id(response)
                        
                        if video_artifact_id:
                            logger.info(
                                "Generated video artifact for Slide %d: %s",
                                slide_idx, video_artifact_id
                            )
                            video_data = video_artifact_id
                        else:
                            logger.warning(
                                "No video artifact in response for Slide %d",
                                slide_idx
                            )
                
                except Exception as e:
                    logger.warning(
                        "MCP video generation failed for Slide %d: %s. "
                        "Saving prompt only.",
                        slide_idx, str(e)
                    )
                
                # Save video prompt for reference
                video_prompt_file = os.path.join(
                    videos_dir,
                    f"slide_{slide_idx}_video_prompt.txt"
                )
                with open(video_prompt_file, "w", encoding="utf-8") as f:
                    f.write("Slide %d Video Prompt\n" % slide_idx)
                    f.write("="*29 + "\n\n")
                    f.write("Prompt:\n%s\n\n" % video_prompt)
                    f.write("Speaker Notes:\n%s\n" % speaker_notes)
                    if video_data:
                        f.write(f"\nGenerated Video: {video_data}\n")
                
                logger.info("Saved video prompt to %s", video_prompt_file)
                
            except Exception:
                logger.error(
                    "Error generating video for Slide %d",
                    slide_idx, exc_info=True
                )
                continue
            finally:
                pass
        
        # Close PDF document
        try:
            pdf_doc.close()
        except Exception:
            pass
        
        logger.info("Video generation phase completed")

    def _extract_video_prompt(
        self, slide_idx: int, speaker_notes: str
    ) -> str:
        """
        Extract a concise video prompt from speaker notes.
        
        Analyzes speaker notes and creates a focused video prompt
        that captures the key visual concepts.
        
        Args:
            slide_idx: Slide index for context
            speaker_notes: Full speaker notes for the slide
            
        Returns:
            Concise video prompt (1-2 sentences)
        """
        return extract_video_prompt(speaker_notes)

    def _extract_artifact_id(self, agent_response: str) -> str:
        """
        Extract artifact_id or video reference from agent response.
        
        Parses agent response to find references to generated video artifacts.
        Looks for patterns like 'artifact_id', 'video_', or common file patterns.
        
        Args:
            agent_response: Full text response from the video agent
            
        Returns:
            Artifact ID/filename if found, empty string otherwise
        """
        return extract_artifact_id(agent_response)
