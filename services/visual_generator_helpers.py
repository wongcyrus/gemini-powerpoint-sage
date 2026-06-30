"""Pure helpers for visual generation prompt and placement logic."""

from __future__ import annotations

import os

from PIL import Image
from config.constants import LanguageConfig


def get_logo_instruction(slide_idx: int) -> str:
    """Return the branding instruction for a slide position."""
    if slide_idx == 1:
        return (
            "You MUST prominently feature the logo/branding from "
            "IMAGE 1 (Original Draft Slide) in an appropriate corner."
        )
    return "DO NOT include any logos or branding elements. Focus solely on content."


def build_designer_prompt(
    speaker_notes: str,
    logo_instruction: str,
    language: str = "en",
    previous_image_present: bool = False,
) -> str:
    """Build the prompt for the primary visual designer agent."""
    style_ref = "Style Reference (Previous Slide) provided." if previous_image_present else "N/A"
    lang_name = LanguageConfig.get_language_name(language)
    lang_instruction = ""
    if language != "en":
        lang_instruction = (
            f"\n\nLANGUAGE: ALL text in the generated image MUST be "
            f"in {lang_name}. Do NOT include any English text. "
            f"Translate all titles, labels, and content to {lang_name}."
        )

    return (
        f"IMAGE 1: Original Slide Image provided.\n\n"
        f"IMAGE 2: {style_ref}\n\n"
        f"Speaker Notes: \"{speaker_notes}\"\n\n"
        f"TASK: Generate the high-fidelity slide image now.\n\n"
        f"CONTEXT: {logo_instruction}{lang_instruction}\n"
    )


def compute_image_placement_inches(
    img_width_px: int,
    img_height_px: int,
    slide_width_inches: float,
    slide_height_inches: float,
    dpi: int = 96,
    mode: str = "cover",
) -> tuple[float, float, float, float]:
    """Compute image placement in inches for PowerPoint."""
    slide_width_px = slide_width_inches * dpi
    slide_height_px = slide_height_inches * dpi

    img_ratio = img_width_px / img_height_px
    slide_ratio = slide_width_px / slide_height_px

    if mode == "contain":
        if img_ratio > slide_ratio:
            new_width_px = slide_width_px
            new_height_px = slide_width_px / img_ratio
            left_px = 0
            top_px = (slide_height_px - new_height_px) / 2
        else:
            new_height_px = slide_height_px
            new_width_px = slide_height_px * img_ratio
            left_px = (slide_width_px - new_width_px) / 2
            top_px = 0
    else:
        if img_ratio > slide_ratio:
            new_height_px = slide_height_px
            new_width_px = slide_height_px * img_ratio
            left_px = -(new_width_px - slide_width_px) / 2
            top_px = 0
        else:
            new_width_px = slide_width_px
            new_height_px = slide_width_px / img_ratio
            left_px = 0
            top_px = -(new_height_px - slide_height_px) / 2

    return (
        left_px / dpi,
        top_px / dpi,
        new_width_px / dpi,
        new_height_px / dpi,
    )


def optimize_image_file(img_path: str) -> str:
    """Re-save an image with smaller file size when possible."""
    reduced_img_path = img_path
    try:
        with Image.open(img_path) as im:
            if im.mode in ("RGBA", "LA") or (im.format == "PNG" and "transparency" in im.info):
                tmp_path = os.path.splitext(img_path)[0] + "_reduced.png"
                im.save(tmp_path, format="PNG", optimize=True)
            else:
                tmp_path = os.path.splitext(img_path)[0] + "_reduced.jpg"
                im = im.convert("RGB")
                im.save(tmp_path, format="JPEG", quality=85, optimize=True)
            reduced_img_path = tmp_path
    except Exception:
        return img_path

    return reduced_img_path


def cleanup_reduced_image_file(original_path: str, reduced_path: str) -> None:
    """Delete the generated reduced image when it differs from the source."""
    if reduced_path != original_path and os.path.exists(reduced_path):
        os.unlink(reduced_path)


def apply_slide_notes(slide, notes: str) -> None:
    """Write plain notes text onto a slide notes section."""
    if not slide.has_notes_slide:
        slide.notes_slide

    text_frame = slide.notes_slide.notes_text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    paragraph.text = notes
    paragraph.level = 0

    from pptx.enum.text import PP_ALIGN

    paragraph.alignment = PP_ALIGN.LEFT
