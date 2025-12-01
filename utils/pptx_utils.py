"""
Utility functions for PowerPoint file manipulation and VBA macro preservation.
"""

import logging
import os
import shutil
import tempfile
import zipfile
import xml.etree.ElementTree as ET
from typing import Optional

from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)

def ensure_pptx_path(path: str) -> str:
    """
    Ensure the path has a .pptx extension.
    
    Args:
        path: The file path to check.
        
    Returns:
        The path with .pptx extension.
    """
    base, ext = os.path.splitext(path)
    if ext.lower() == ".pptx":
        return path
    return base + ".pptx"

def get_slide_notes(slide) -> str:
    """
    Extract existing notes from a slide.
    
    Args:
        slide: The slide object.
        
    Returns:
        The extracted notes text, or empty string if none found.
    """
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        return slide.notes_slide.notes_text_frame.text.strip()
    return ""

def update_slide_notes(slide, notes: str) -> None:
    """
    Update the notes on a slide as plain text without bullets.
    
    Args:
        slide: The slide object to update.
        notes: The new notes text.
    """
    try:
        if not slide.has_notes_slide:
            slide.notes_slide
        
        text_frame = slide.notes_slide.notes_text_frame
        text_frame.clear()
        
        # Add notes as single paragraph without bullet formatting
        p = text_frame.paragraphs[0]
        p.text = notes
        p.level = 0
        
        # Explicitly remove bullet formatting
        p.alignment = PP_ALIGN.LEFT
        
    except Exception as e:
        logger.error(f"Failed to update slide notes: {e}")

def restore_vba_project(original_src: str, generated_pptx: str, final_path: str) -> None:
    """
    Inject VBA macros (vbaProject.bin) from original source into the generated PPTX.
    
    Args:
        original_src: Path to the original PPTM file containing macros.
        generated_pptx: Path to the newly generated PPTX file.
        final_path: Path where the final PPTM should be saved.
    """
    # If original source contains vbaProject.bin, inject it and write final_path.
    try:
        with zipfile.ZipFile(original_src, 'r') as zorig:
            try:
                vba_data = zorig.read('ppt/vbaProject.bin')
            except KeyError:
                vba_data = None

        if not vba_data:
            # No macros present; just move/rename generated file to final path
            # Ensure we don't overwrite if source and dest are same (unlikely here but good practice)
            if generated_pptx != final_path:
                shutil.move(generated_pptx, final_path)
            return

        # Extract generated pptx to temp dir
        tmpdir = tempfile.mkdtemp(prefix='pptm_merge_')
        try:
            with zipfile.ZipFile(generated_pptx, 'r') as zgen:
                zgen.extractall(tmpdir)

            # Ensure ppt/ exists and write vbaProject.bin
            ppt_dir = os.path.join(tmpdir, 'ppt')
            os.makedirs(ppt_dir, exist_ok=True)
            with open(os.path.join(ppt_dir, 'vbaProject.bin'), 'wb') as f:
                f.write(vba_data)

            # Update [Content_Types].xml
            ctype_path = os.path.join(tmpdir, '[Content_Types].xml')
            if os.path.exists(ctype_path):
                tree = ET.parse(ctype_path)
                root = tree.getroot()
                # Namespace handling
                override_tag = '{http://schemas.openxmlformats.org/package/2006/content-types}Override'
                exists = False
                for ov in root.findall(override_tag):
                    if ov.attrib.get('PartName') == '/ppt/vbaProject.bin':
                        exists = True
                        break
                if not exists:
                    ET.register_namespace('', 'http://schemas.openxmlformats.org/package/2006/content-types')
                    new = ET.Element(override_tag, PartName='/ppt/vbaProject.bin', ContentType='application/vnd.ms-office.vbaProject')
                    root.append(new)
                    tree.write(ctype_path, encoding='utf-8', xml_declaration=True)

            # Update ppt/_rels/presentation.xml.rels
            rels_dir = os.path.join(tmpdir, 'ppt', '_rels')
            os.makedirs(rels_dir, exist_ok=True)
            pres_rels = os.path.join(rels_dir, 'presentation.xml.rels')
            rel_tag = '{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'

            if os.path.exists(pres_rels):
                tree = ET.parse(pres_rels)
                root = tree.getroot()
            else:
                # Create minimal rels structure
                root = ET.Element('{http://schemas.openxmlformats.org/package/2006/relationships}Relationships')
                tree = ET.ElementTree(root)

            # Check if vba relationship already exists
            exists = False
            for r in root.findall(rel_tag):
                if r.attrib.get('Type') == 'http://schemas.microsoft.com/office/2006/relationships/vbaProject':
                    exists = True
                    break

            if not exists:
                # generate a relationship Id
                rid = 'rIdVBA1'
                newr = ET.Element(rel_tag, Id=rid, Type='http://schemas.microsoft.com/office/2006/relationships/vbaProject', Target='vbaProject.bin')
                root.append(newr)
                tree.write(pres_rels, encoding='utf-8', xml_declaration=True)

            # Ensure parent directory exists before creating zipfile
            final_dir = os.path.dirname(final_path)
            if final_dir:
                os.makedirs(final_dir, exist_ok=True)

            # Repack into final_path
            with zipfile.ZipFile(final_path, 'w', zipfile.ZIP_DEFLATED) as zout:
                for foldername, subfolders, filenames in os.walk(tmpdir):
                    for filename in filenames:
                        file_path = os.path.join(foldername, filename)
                        arcname = os.path.relpath(file_path, tmpdir)
                        zout.write(file_path, arcname)

        finally:
            # Cleanup
            shutil.rmtree(tmpdir)
            try:
                if os.path.exists(generated_pptx):
                    os.remove(generated_pptx)
            except Exception:
                pass

    except Exception as e:
        logger.error('Failed to preserve VBA project: %s', e, exc_info=True)
        # Fallback: just move the generated file
        try:
            if generated_pptx != final_path:
                shutil.move(generated_pptx, final_path)
        except Exception:
            pass
